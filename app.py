# app.py - single-file Streamlit app embedding the extraction pipeline

import os
import io
import tempfile
import zipfile
import statistics
from zipfile import ZipFile
from urllib.parse import urlparse, urljoin
from pathlib import Path

import streamlit as st
from dotenv import load_dotenv
load_dotenv()

import requests
from bs4 import BeautifulSoup
from PIL import Image
import fitz  # PyMuPDF
import docx
from pptx import Presentation

# Google GenAI client
from google import genai
from google.genai import types

# CONFIG
MODEL = "gemini-2.0-flash"
GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY")
if GOOGLE_API_KEY:
    client = genai.Client(api_key=GOOGLE_API_KEY)
else:
    client = None  # used only when not in dry-run

# ------------------------
# Helper / pipeline functions
# ------------------------

def _download_url_to_tempfile(url, timeout=10):
    """Download a URL and save to a temp file. Return (tmp_path, ext)."""
    try:
        head = requests.head(url, allow_redirects=True, timeout=5)
    except Exception:
        head = None

    content_type = None
    if head is not None and 'content-type' in head.headers:
        content_type = head.headers['content-type']

    if not content_type:
        resp = requests.get(url, stream=True, timeout=timeout)
    else:
        resp = requests.get(url, stream=True, timeout=timeout)

    if resp.status_code >= 400:
        raise ValueError(f"Failed to fetch URL: {url} (status {resp.status_code})")

    ct = resp.headers.get('content-type', content_type or '').lower()

    if 'html' in ct:
        ext = '.html'
    elif 'pdf' in ct:
        ext = '.pdf'
    elif 'presentation' in ct or 'pptx' in ct:
        ext = '.pptx'
    elif 'word' in ct or 'officedocument.wordprocessingml.document' in ct or 'msword' in ct:
        ext = '.docx'
    elif 'image' in ct:
        subtype = ct.split('/')[-1].split(';')[0] if '/' in ct else 'jpeg'
        ext = f".{subtype}" if not subtype.startswith('x-') else f".{subtype.replace('x-', '')}"
    else:
        path = urlparse(url).path
        _, ext = os.path.splitext(path)
        if not ext:
            ext = '.html'

    tmp_fd, tmp_path = tempfile.mkstemp(suffix=ext)
    os.close(tmp_fd)

    with open(tmp_path, 'wb') as out_f:
        for chunk in resp.iter_content(chunk_size=8192):
            if chunk:
                out_f.write(chunk)

    return tmp_path, ext


# --- extractors (min_size required) ---

def extract_figures_from_pdf(pdf_path, output_dir, min_size):
    if not isinstance(min_size, int) or min_size < 0:
        raise ValueError("min_size must be a non-negative integer")
    os.makedirs(output_dir, exist_ok=True)
    doc = fitz.open(pdf_path)
    image_data = []
    skipped = 0
    for page_index, page in enumerate(doc):
        images = page.get_images(full=True)
        for img_index, img in enumerate(images):
            xref = img[0]
            smask = img[1]
            try:
                base_pix = fitz.Pixmap(doc, xref)
            except Exception:
                continue
            if getattr(base_pix, "alpha", False):
                try:
                    base_pix = fitz.Pixmap(base_pix, 0)
                except Exception:
                    pass
            pix = base_pix
            if smask and smask > 0:
                try:
                    mask_pix = fitz.Pixmap(doc, smask)
                    try:
                        pix = fitz.Pixmap(base_pix, mask_pix)
                    except Exception:
                        pix = base_pix
                    finally:
                        mask_pix = None
                except Exception:
                    pix = base_pix
            try:
                if getattr(pix, "colorspace", None) and pix.colorspace.n > 3:
                    pix = fitz.Pixmap(fitz.csRGB, pix)
            except Exception:
                pass
            try:
                sample_len = len(pix.samples)
            except Exception:
                sample_len = 0
            if sample_len < min_size:
                skipped += 1
                pix = None
                base_pix = None
                continue
            temp_path = os.path.join(output_dir, f"temp_page_{page_index+1}_figure_{img_index+1}.png")
            try:
                pix.save(temp_path)
            except Exception:
                pix = None
                base_pix = None
                try:
                    if os.path.exists(temp_path):
                        os.remove(temp_path)
                except Exception:
                    pass
                continue
            try:
                pil_img = Image.open(temp_path)
            except Exception:
                try:
                    os.remove(temp_path)
                except Exception:
                    pass
                pix = None
                base_pix = None
                continue
            if pil_img.mode in ('RGBA', 'LA', 'P'):
                if pil_img.mode == 'P':
                    pil_img = pil_img.convert('RGBA')
                background = Image.new("RGB", pil_img.size, (255, 255, 255))
                try:
                    alpha = pil_img.split()[3] if pil_img.mode == 'RGBA' else pil_img.split()[1]
                    background.paste(pil_img, mask=alpha)
                except Exception:
                    background.paste(pil_img)
                pil_img = background
            img_ext = "jpeg"
            img_path = os.path.join(output_dir, f"page_{page_index+1}_figure_{img_index+1}.{img_ext}")
            try:
                pil_img.save(img_path, "JPEG", quality=95)
            except Exception:
                try:
                    os.remove(temp_path)
                except Exception:
                    pass
                pix = None
                base_pix = None
                continue
            try:
                with open(img_path, "rb") as f:
                    img_bytes = f.read()
            except Exception:
                img_bytes = b''
            image_data.append((img_path, img_bytes, img_ext, page_index))
            try:
                os.remove(temp_path)
            except Exception:
                pass
            pix = None
            base_pix = None
    doc.close()
    return image_data


def extract_figures_from_docx(docx_path, output_dir, min_size):
    if not isinstance(min_size, int) or min_size < 0:
        raise ValueError("min_size must be a non-negative integer")
    os.makedirs(output_dir, exist_ok=True)
    image_data = []
    try:
        with ZipFile(docx_path) as zipf:
            media_files = [name for name in zipf.namelist() if name.startswith('word/media/')]
            for img_index, media in enumerate(media_files, 1):
                try:
                    img_bytes = zipf.read(media)
                except Exception:
                    continue
                if len(img_bytes) < min_size:
                    continue
                try:
                    pil_img = Image.open(io.BytesIO(img_bytes))
                except Exception:
                    continue
                if pil_img.mode in ('RGBA', 'LA', 'P'):
                    if pil_img.mode == 'P':
                        pil_img = pil_img.convert('RGBA')
                    background = Image.new("RGB", pil_img.size, (255, 255, 255))
                    try:
                        alpha = pil_img.split()[3] if pil_img.mode == 'RGBA' else pil_img.split()[1]
                        background.paste(pil_img, mask=alpha)
                    except Exception:
                        background.paste(pil_img)
                    pil_img = background
                img_ext = "jpeg"
                base_name = os.path.splitext(os.path.basename(media))[0]
                img_path = os.path.join(output_dir, f"docx_{base_name}_{img_index}.{img_ext}")
                try:
                    pil_img.save(img_path, "JPEG", quality=95)
                except Exception:
                    continue
                try:
                    with open(img_path, "rb") as f:
                        saved_bytes = f.read()
                except Exception:
                    saved_bytes = img_bytes
                image_data.append((img_path, saved_bytes, img_ext, 0))
    except Exception as e:
        st.warning(f"DOCX extraction error: {e}")
    return image_data


def extract_figures_from_pptx(pptx_path, output_dir, min_size):
    if not isinstance(min_size, int) or min_size < 0:
        raise ValueError("min_size must be a non-negative integer")
    os.makedirs(output_dir, exist_ok=True)
    image_data = []
    try:
        with ZipFile(pptx_path) as zipf:
            media_files = [name for name in zipf.namelist() if name.startswith('ppt/media/')]
            for img_index, media in enumerate(media_files, 1):
                try:
                    img_bytes = zipf.read(media)
                except Exception:
                    continue
                if len(img_bytes) < min_size:
                    continue
                try:
                    pil_img = Image.open(io.BytesIO(img_bytes))
                except Exception:
                    continue
                if pil_img.mode in ('RGBA', 'LA', 'P'):
                    if pil_img.mode == 'P':
                        pil_img = pil_img.convert('RGBA')
                    background = Image.new("RGB", pil_img.size, (255, 255, 255))
                    try:
                        alpha = pil_img.split()[3] if pil_img.mode == 'RGBA' else pil_img.split()[1]
                        background.paste(pil_img, mask=alpha)
                    except Exception:
                        background.paste(pil_img)
                    pil_img = background
                img_ext = "jpeg"
                base_name = os.path.splitext(os.path.basename(media))[0]
                img_path = os.path.join(output_dir, f"pptx_{base_name}_{img_index}.{img_ext}")
                try:
                    pil_img.save(img_path, "JPEG", quality=95)
                except Exception:
                    continue
                try:
                    with open(img_path, "rb") as f:
                        saved_bytes = f.read()
                except Exception:
                    saved_bytes = img_bytes
                image_data.append((img_path, saved_bytes, img_ext, 0))
    except Exception as e:
        st.warning(f"PPTX extraction error: {e}")
    return image_data


def extract_figures_from_html(html_path, output_dir, min_size):
    if not isinstance(min_size, int) or min_size < 0:
        raise ValueError("min_size must be a non-negative integer")
    os.makedirs(output_dir, exist_ok=True)
    with open(html_path, 'r', encoding='utf-8') as f:
        soup = BeautifulSoup(f, 'html.parser')
    image_data = []
    base_dir = os.path.dirname(html_path)
    img_index = 0
    for img_tag in soup.find_all('img'):
        src = img_tag.get('src') or img_tag.get('data-src') or img_tag.get('data-original')
        if not src:
            continue
        if src.startswith('//'):
            src = 'https:' + src
        img_bytes = None
        if src.startswith('http://') or src.startswith('https://'):
            try:
                resp = requests.get(src, timeout=5)
                if resp.status_code == 200:
                    img_bytes = resp.content
            except Exception:
                img_bytes = None
        else:
            local_path = os.path.join(base_dir, src)
            if os.path.exists(local_path):
                try:
                    with open(local_path, 'rb') as lf:
                        img_bytes = lf.read()
                except Exception:
                    img_bytes = None
        if not img_bytes:
            continue
        if len(img_bytes) < min_size:
            continue
        try:
            pil_img = Image.open(io.BytesIO(img_bytes))
        except Exception:
            continue
        if pil_img.mode in ('RGBA', 'LA', 'P'):
            if pil_img.mode == 'P':
                pil_img = pil_img.convert('RGBA')
            background = Image.new("RGB", pil_img.size, (255, 255, 255))
            try:
                alpha = pil_img.split()[3] if pil_img.mode == 'RGBA' else pil_img.split()[1]
                background.paste(pil_img, mask=alpha)
            except Exception:
                background.paste(pil_img)
            pil_img = background
        img_ext = "jpeg"
        img_index += 1
        img_path = os.path.join(output_dir, f"html_image_{img_index}.{img_ext}")
        try:
            pil_img.save(img_path, "JPEG", quality=95)
        except Exception:
            continue
        try:
            with open(img_path, "rb") as f:
                saved_bytes = f.read()
        except Exception:
            saved_bytes = img_bytes
        image_data.append((img_path, saved_bytes, img_ext, 0))
    return image_data


# --- context and Gemini call ---

def get_context(doc, context_id, ext):
    if ext == '.pdf':
        page = doc[context_id]
        return page.get_text()
    elif ext == '.docx':
        return '\n'.join([p.text for p in doc.paragraphs])
    elif ext == '.pptx':
        text = ''
        for slide in doc.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + '\n'
        return text
    elif ext == '.html':
        # doc is BeautifulSoup here; return visible text
        return doc.get_text(separator='\n')
    else:
        return ''


def gemini_describe(img_bytes, img_ext, page_text, model=MODEL):
    if client is None:
        raise RuntimeError("GOOGLE_API_KEY not configured (Gemini unavailable).")
    mime_type = f"image/{img_ext}"
    prompt = f"""You are an expert document analyst.
Analyze the image and extract any visible text (labels, legends, annotations). Then describe the image and explain how it relates to the surrounding document text.

Surrounding text:
{page_text}

Return as plain text description."""
    response = client.models.generate_content(
        model=model,
        contents=[
            types.Part(text=prompt),
            types.Part(inline_data=types.Blob(mime_type=mime_type, data=img_bytes))
        ]
    )
    return response.text.strip()


# --- dynamic min_size ---

def determine_min_size(file_path, ext, mode='balanced', sample_limit=8, min_absolute=5000):
    sizes = []
    try:
        if ext == '.pdf':
            doc = fitz.open(file_path)
            for page in doc:
                if len(sizes) >= sample_limit:
                    break
                for img in page.get_images(full=True):
                    if len(sizes) >= sample_limit:
                        break
                    xref = img[0]
                    try:
                        pix = fitz.Pixmap(doc, xref)
                        sizes.append(len(pix.samples))
                    except Exception:
                        continue
            doc.close()
        elif ext == '.docx':
            try:
                doc = docx.Document(file_path)
                for rel in list(doc.part.rels.values()):
                    if len(sizes) >= sample_limit:
                        break
                    if "image" in rel.reltype:
                        try:
                            blob = rel.target_part.blob
                            sizes.append(len(blob))
                        except Exception:
                            continue
            except Exception:
                pass
        elif ext == '.pptx':
            try:
                with ZipFile(file_path) as zipf:
                    media_files = [name for name in zipf.namelist() if name.startswith('ppt/media/')]
                    for media in media_files[:sample_limit]:
                        try:
                            b = zipf.read(media)
                            sizes.append(len(b))
                        except Exception:
                            continue
            except Exception:
                pass
        elif ext in ('.html', '.htm'):
            try:
                with open(file_path, 'r', encoding='utf-8') as f:
                    soup = BeautifulSoup(f, 'html.parser')
                img_tags = soup.find_all('img')[:sample_limit]
                for tag in img_tags:
                    src = (tag.get('src') or tag.get('data-src') or tag.get('data-original') or '').strip()
                    if not src:
                        continue
                    if src.startswith('//'):
                        src = 'https:' + src
                    if src.startswith('http'):
                        try:
                            head = requests.head(src, allow_redirects=True, timeout=3)
                            cl = head.headers.get('content-length')
                            if cl and cl.isdigit():
                                sizes.append(int(cl))
                                continue
                        except Exception:
                            pass
                        try:
                            resp = requests.get(src, stream=True, timeout=5)
                            if resp.status_code == 200:
                                sizes.append(len(resp.content))
                        except Exception:
                            continue
                    else:
                        local_path = os.path.join(os.path.dirname(file_path), src)
                        if os.path.exists(local_path):
                            try:
                                with open(local_path, 'rb') as lf:
                                    sizes.append(len(lf.read()))
                            except Exception:
                                continue
            except Exception:
                pass
    except Exception:
        sizes = []
    if mode == 'aggressive':
        multiplier = 0.25
    elif mode == 'conservative':
        multiplier = 1.2
    else:
        multiplier = 0.5
    if sizes:
        med = int(statistics.median(sizes))
        computed = max(int(med * multiplier), min_absolute)
        computed = int(max(min_absolute, min(computed, max(200000, med * 4))))
        return computed
    else:
        fallback_map = {'.pdf': 50000, '.docx': 15000, '.pptx': 20000, '.html': 8000}
        return fallback_map.get(ext, min_absolute)


# --- process_document uses determine_min_size and extractors ---

def process_document(file_path, determine_mode='balanced', dry_run=False):
    temp_file_to_remove = None
    is_url = isinstance(file_path, str) and file_path.lower().startswith(('http://', 'https://'))
    original_url = file_path if is_url else None
    if is_url:
        tmp_path, tmp_ext = _download_url_to_tempfile(file_path)
        temp_file_to_remove = tmp_path
        file_path_to_use = tmp_path
    else:
        file_path_to_use = file_path
    try:
        ext = os.path.splitext(file_path_to_use)[1].lower()
        base_output_dir = f"extracted_images_{ext[1:] if ext else 'unknown'}"
        if is_url:
            parsed = urlparse(file_path)
            url_basename = os.path.basename(parsed.path) or parsed.netloc
            file_folder = os.path.splitext(url_basename)[0]
            file_folder = "".join(c for c in file_folder if c.isalnum() or c in ('_', '-')).strip() or "downloaded_file"
        else:
            file_folder = os.path.splitext(os.path.basename(file_path_to_use))[0]
        output_dir = os.path.join(base_output_dir, file_folder)
        os.makedirs(output_dir, exist_ok=True)
        output_file = os.path.join(output_dir, f"extracted_texts_{file_folder}.txt")
        min_size = determine_min_size(file_path_to_use, ext, mode=determine_mode, sample_limit=8, min_absolute=5000)
        image_data = []
        if ext == '.pdf':
            image_data = extract_figures_from_pdf(file_path_to_use, output_dir, min_size=min_size)
            doc = fitz.open(file_path_to_use)
        elif ext == '.docx':
            image_data = extract_figures_from_docx(file_path_to_use, output_dir, min_size=min_size)
            doc = docx.Document(file_path_to_use)
        elif ext == '.pptx':
            image_data = extract_figures_from_pptx(file_path_to_use, output_dir, min_size=min_size)
            doc = Presentation(file_path_to_use)
        elif ext in ('.html', '.htm'):
            image_data = extract_figures_from_html(file_path_to_use, output_dir, min_size=min_size)
            with open(file_path_to_use, 'r', encoding='utf-8') as f:
                doc = BeautifulSoup(f, 'html.parser')
        else:
            try:
                with open(file_path_to_use, 'r', encoding='utf-8') as f:
                    _ = f.read(2048)
                ext = '.html'
                image_data = extract_figures_from_html(file_path_to_use, output_dir, min_size=min_size)
                with open(file_path_to_use, 'r', encoding='utf-8') as f:
                    doc = BeautifulSoup(f, 'html.parser')
            except Exception:
                raise ValueError(f"Unsupported or unrecognized file type for: {file_path}")
        descriptions = []
        if dry_run:
            for img_path, img_bytes, img_ext, context_id in image_data:
                descriptions.append((img_path, "(dry-run)"))
        else:
            if client is None:
                raise RuntimeError("GOOGLE_API_KEY not configured; enable dry-run or set API key.")
            for img_path, img_bytes, img_ext, context_id in image_data:
                context_text = get_context(doc, context_id, ext)
                desc = gemini_describe(img_bytes, img_ext, context_text)
                descriptions.append((img_path, desc))
        # Write descriptions to file
        with open(output_file, "w", encoding="utf-8") as out:
            for img_path, desc in descriptions:
                out.write(f"{img_path}:\n{desc}\n\n")
        if ext == '.pdf':
            doc.close()
        return output_dir, output_file, image_data, descriptions
    finally:
        if temp_file_to_remove and os.path.exists(temp_file_to_remove):
            try:
                os.remove(temp_file_to_remove)
            except Exception:
                pass


# ------------------------
# Streamlit UI
# ------------------------

st.set_page_config(page_title="Image Extraction Pipeline", layout="wide")
st.title("Image Extraction + Gemini (single-file)")

with st.sidebar:
    st.header("Options")
    mode = st.selectbox("Extraction mode", ["balanced", "aggressive", "conservative"])
    dry_run = st.checkbox("Dry run (skip Gemini calls)", value=True)
    keep_temp = st.checkbox("Keep downloaded temp files (debug)", value=False)

uploaded_file = st.file_uploader("Upload a file (PDF, DOCX, PPTX, HTML)", type=["pdf", "docx", "pptx", "html", "htm"])
url_input = st.text_input("Or enter a URL to process (http/https)")
run = st.button("Run Extraction")

def _save_uploaded_tmp(uploaded):
    tmp_fd, tmp_path = tempfile.mkstemp(suffix="_" + uploaded.name)
    os.close(tmp_fd)
    with open(tmp_path, "wb") as f:
        f.write(uploaded.getbuffer())
    return tmp_path

if run:
    tmp_files = []
    try:
        if uploaded_file is None and not url_input.strip():
            st.error("Provide an upload or a URL.")
        else:
            if uploaded_file:
                st.info(f"Processing uploaded file: {uploaded_file.name}")
                source = _save_uploaded_tmp(uploaded_file)
                tmp_files.append(source)
            else:
                st.info(f"Processing URL: {url_input}")
                source = url_input.strip()
            with st.spinner("Running pipeline..."):
                try:
                    out_dir, out_txt, image_data, descriptions = process_document(source, determine_mode=mode, dry_run=dry_run)
                except Exception as e:
                    st.error(f"Pipeline error: {e}")
                    raise
            st.success("Done.")
            st.write(f"Output directory: `{out_dir}`")
            # show images
            imgs = sorted([f for f in os.listdir(out_dir) if f.lower().endswith((".jpg", ".jpeg", ".png"))])
            col1, col2 = st.columns([1,1])
            with col1:
                st.subheader("Extracted images")
                if not imgs:
                    st.info("No images found or all skipped by threshold.")
                else:
                    for name in imgs:
                        st.image(os.path.join(out_dir, name), caption=name)
            with col2:
                st.subheader("Descriptions")
                st.text_area("Gemini output", value=open(out_txt, "r", encoding="utf-8").read(), height=600)
            # download zip
            zip_path = os.path.join(tempfile.gettempdir(), f"{Path(out_dir).name}.zip")
            with zipfile.ZipFile(zip_path, "w", zipfile.ZIP_DEFLATED) as z:
                for root, _, files in os.walk(out_dir):
                    for fn in files:
                        full = os.path.join(root, fn)
                        rel = os.path.relpath(full, start=out_dir)
                        z.write(full, arcname=rel)
            with open(zip_path, "rb") as f:
                st.download_button("Download results (zip)", data=f.read(), file_name=os.path.basename(zip_path))
    finally:
        if not keep_temp:
            for t in tmp_files:
                try:
                    os.remove(t)
                except Exception:
                    pass

st.markdown("---")
st.markdown("**Notes:** set `GOOGLE_API_KEY` in the environment to enable Gemini calls. Use Dry run to test extraction without API usage.")
